"""文档解析（4.3）：文本/表格提取为 doc_block（M3 起由任务队列调度）。

旧版 Office 格式（.doc/.xls/.ppt）统一经 LibreOffice 无头转新版（docx/xlsx/pptx）后解析，
表格结构保留；转换失败给出可操作提示而非"不支持的格式"。"""

from __future__ import annotations

from contextlib import contextmanager
from pathlib import Path


def parse_to_blocks(path: Path, ext: str) -> list[dict]:
    """返回 [{block_type, page_no, block_index, text_content, extra}]。"""
    ext = ext.lower()
    if ext in (".txt", ".csv"):
        # 中文常见 GBK 编码：utf-8 失败回退 gbk，避免乱码块
        raw = Path(path).read_bytes()
        try:
            text = raw.decode("utf-8")
        except UnicodeDecodeError:
            text = raw.decode("gbk", errors="replace")
        return [{"block_type": "paragraph", "page_no": None, "block_index": 0, "text_content": text}]

    if ext == ".ofd":
        return _parse_ofd(path)

    if ext == ".docx":
        return _parse_docx(path)

    if ext == ".xlsx":
        return _parse_xlsx(path)

    if ext == ".xls":
        with _converted(path, "xlsx") as out:
            blocks = _parse_xlsx(out)
        if not blocks:
            raise ValueError("表格无可提取内容（可能为空白文件或扫描件）")
        for b in blocks:
            b.setdefault("extra", {})["source"] = "libreoffice-xls"
        return blocks

    if ext == ".pdf":
        try:
            import fitz  # pymupdf
        except ImportError as exc:  # pragma: no cover - 容器环境安装
            raise ValueError("PDF 解析库未安装（pymupdf）") from exc
        doc = fitz.open(str(path))
        blocks = []
        idx = 0
        for page_no, page in enumerate(doc, start=1):
            text = page.get_text().strip()
            if text:
                blocks.append(
                    {
                        "block_type": "paragraph",
                        "page_no": page_no,
                        "block_index": idx,
                        "text_content": text,
                    }
                )
                idx += 1
        doc.close()
        return blocks

    if ext == ".pptx":
        return _parse_pptx(path)

    if ext == ".ppt":
        with _converted(path, "pptx") as out:
            blocks = _parse_pptx(out)
        if not blocks:
            raise ValueError("演示文稿无可提取文本（可能为空白文件或纯图片）")
        for b in blocks:
            b.setdefault("extra", {})["source"] = "libreoffice-ppt"
        return blocks

    if ext in (".html", ".htm"):
        # 公告网页导入（Issue #4/#6）：去标签提取正文（纯标准库，无第三方依赖）
        from html.parser import HTMLParser

        class _TextExtractor(HTMLParser):
            def __init__(self) -> None:
                super().__init__()
                self.parts: list[str] = []

            def handle_data(self, data: str) -> None:
                if data.strip():
                    self.parts.append(data.strip())

        p = _TextExtractor()
        p.feed(Path(path).read_text(encoding="utf-8", errors="replace"))
        text = "\n".join(p.parts)
        if not text.strip():
            raise ValueError("HTML 无可提取文本")
        return [{"block_type": "paragraph", "page_no": None, "block_index": 0, "text_content": text, "extra": {"source": "html-text"}}]

    if ext == ".doc":
        return _parse_legacy_doc(path)

    raise ValueError(f"不支持的格式：{ext}")


@contextmanager
def _converted(path: Path, dst_fmt: str):
    """LibreOffice 无头转换到临时目录，产出 source.<dst_fmt>；退出自动清理。

    无 LibreOffice 或转换失败时给出可操作提示（此前旧格式直接报"不支持的格式"，
    用户无从下手——Issue #8 任务 224 同类问题）。"""
    import shutil
    import subprocess
    import tempfile

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice is None:
        raise ValueError("旧版 Office 文件需 LibreOffice 转换，但服务器未安装转换组件；请用 Office 另存为新格式后重新上传")
    with tempfile.TemporaryDirectory() as tmp:
        src = Path(tmp) / f"source{path.suffix.lower() or ''}"
        src.write_bytes(path.read_bytes())
        try:
            proc = subprocess.run(
                [soffice, "--headless", "--norestore", "--convert-to", dst_fmt, "--outdir", tmp, str(src)],
                capture_output=True,
                timeout=180,
            )
        except subprocess.TimeoutExpired as exc:
            raise ValueError("旧版 Office 文件转换超时；请用 Office 另存为新格式后重新上传") from exc
        if proc.returncode != 0:
            raise ValueError("旧版 Office 文件转换失败；请用 Office 另存为新格式后重新上传")
        out = Path(tmp) / f"source.{dst_fmt}"
        if not out.exists():
            raise ValueError("旧版 Office 文件转换未产出目标格式；请用 Office 另存为新格式后重新上传")
        yield out


def _parse_xlsx(path: Path) -> list[dict]:
    """xlsx：各 sheet 行提取为 table block。"""
    import io as _io

    from openpyxl import load_workbook

    # 生产存储对象路径为内容寻址的 ".../<sha256>/original"（无扩展名）；
    # openpyxl 会按“文件名扩展名”做格式校验并拒绝无扩展名路径。
    # 从字节流加载可绕过文件名检查；本地测试用带扩展名临时文件，永远踩不到此差异。
    wb = load_workbook(_io.BytesIO(Path(path).read_bytes()), read_only=True, data_only=True)
    blocks = []
    idx = 0
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            vals = ["" if v is None else str(v) for v in row]
            if any(vals):
                blocks.append(
                    {
                        "block_type": "table",
                        "page_no": None,
                        "block_index": idx,
                        "text_content": " | ".join(vals),
                        "extra": {"sheet": ws.title, "cols": vals},
                    }
                )
                idx += 1
    wb.close()
    return blocks


def _parse_pptx(path: Path) -> list[dict]:
    """pptx：每页文本+表格提取为块；python-pptx 不可用时纯标准库兜底。"""
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        blocks = []
        idx = 0
        for slide_no, slide in enumerate(prs.slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    parts.append(shape.text_frame.text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        parts.append(" | ".join(cell.text for cell in row.cells))
            text = "\n".join(p for p in parts if p.strip())
            if text.strip():
                blocks.append(
                    {
                        "block_type": "paragraph",
                        "page_no": slide_no,
                        "block_index": idx,
                        "text_content": text,
                        "extra": {"source": "pptx-slide", "page": slide_no},
                    }
                )
                idx += 1
        if blocks:
            return blocks
    except Exception:  # noqa: BLE001 库缺失或包结构异常 → 纯标准库兜底
        pass
    return _parse_pptx_stdlib(path)


def _parse_docx(path: Path) -> list[dict]:
    """docx：段落 + 表格行提取为 doc_block。"""
    from docx import Document

    doc = Document(str(path))
    blocks: list[dict] = []
    idx = 0
    for para in doc.paragraphs:
        if para.text.strip():
            blocks.append(
                {"block_type": "paragraph", "page_no": None, "block_index": idx, "text_content": para.text}
            )
            idx += 1
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            blocks.append(
                {
                    "block_type": "table",
                    "page_no": None,
                    "block_index": idx,
                    "text_content": " | ".join(cells),
                    "extra": {"cols": cells},
                }
            )
            idx += 1
    # 页眉/页脚（正文解析不覆盖，作为索引补充块）
    for section in doc.sections:
        for part_name, part in (("header", section.header), ("footer", section.footer)):
            if part is None or part.is_linked_to_previous:
                continue
            for para in part.paragraphs:
                if para.text.strip():
                    blocks.append(
                        {
                            "block_type": "header_footer",
                            "page_no": None,
                            "block_index": idx,
                            "text_content": para.text,
                            "extra": {"source": f"docx-{part_name}"},
                        }
                    )
                    idx += 1
    return blocks


def _parse_legacy_doc(path: Path) -> list[dict]:
    """旧版 Word .doc（OLE2 二进制）：LibreOffice 无头转 .docx 后按 docx 解析
    （表格结构保留；此前转 TXT 会丢表格）。转换失败给出可操作提示
    （Issue #8 任务 224：合同条款（空白）.doc 曾报"不支持的格式：.doc"）。"""
    with _converted(path, "docx") as out:
        blocks = _parse_docx(out)
    if not blocks:
        raise ValueError("文档无可提取文本（可能为空白文档或扫描件）；如为扫描件请走图片解析")
    for b in blocks:
        b.setdefault("extra", {})["source"] = "libreoffice-doc"
    return blocks


def _parse_pptx_stdlib(path: Path) -> list[dict]:
    """pptx 文本层提取（纯标准库）：解 zip，按 slideN.xml 顺序提取 <a:t> 文本。

    无需 python-pptx 依赖即可支持 .pptx（生产 Issue #13：供应商注意事项 .pptx 上传后
    解析报"不支持的格式"，任务失败且原因不可见）。
    """
    import re
    import xml.etree.ElementTree as ET
    import zipfile

    def local(tag: str) -> str:
        return tag.rsplit("}", 1)[-1]

    blocks: list[dict] = []
    idx = 0
    with zipfile.ZipFile(path) as zf:
        slide_names = sorted(
            (n for n in zf.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)),
            key=lambda n: int(re.search(r"slide(\d+)\.xml$", n).group(1)),
        )
        for name in slide_names:
            root = ET.fromstring(zf.read(name))
            texts = [
                ("".join(t.itertext()) or "").strip()
                for t in root.iter()
                if local(t.tag) == "t"
            ]
            texts = [t for t in texts if t]
            if not texts:
                continue
            page_no = int(re.search(r"slide(\d+)\.xml$", name).group(1))
            blocks.append(
                {
                    "block_type": "paragraph",
                    "page_no": page_no,
                    "block_index": idx,
                    "text_content": "\n".join(texts),
                    "extra": {"source": "pptx-text-layer", "page": page_no},
                }
            )
            idx += 1
    if not blocks:
        raise ValueError("pptx 无文本层（图片型幻灯片请走视觉模型）")
    return blocks


def _parse_ofd(path: Path) -> list[dict]:
    """OFD 文本层提取（纯标准库，zip+XML）。

    OFD 版式规范：OFD.xml -> Doc_0/Document.xml -> Pages/Page_N/Content.xml，
    TextObject/TextCode 承载文本。命名空间各实现不一（标准 http://www.ofdspec.org/2016、
    easyofd 自生成 http://blog.yuanhaiying.cn 等），按本地名匹配。
    """
    import xml.etree.ElementTree as ET
    import zipfile

    def local(tag: str) -> str:
        return tag.rsplit("}", 1)[-1]

    def elem_text(e: ET.Element) -> str:
        return "".join(e.itertext()).strip()

    blocks: list[dict] = []
    idx = 0
    with zipfile.ZipFile(path) as zf:
        names = zf.namelist()
        doc_root = None
        if "OFD.xml" in names:
            root = ET.fromstring(zf.read("OFD.xml"))
            for el in root.iter():
                if local(el.tag) == "DocRoot":
                    doc_root = (elem_text(el) or "").strip().lstrip("/") or None
                    if doc_root:
                        break
        if not doc_root:
            doc_root = next((n for n in names if n.endswith("Document.xml")), None)
        if not doc_root:
            raise ValueError("OFD 缺少 Document.xml")

        doc_xml = zf.read(doc_root)
        doc = ET.fromstring(doc_xml)
        page_locs: list[str] = []
        tpl_locs: list[str] = []
        for el in doc.iter():
            name = local(el.tag)
            if name == "Page":
                loc = (el.get("BaseLoc") or elem_text(el) or "").strip().lstrip("/")
                if loc and loc not in page_locs:
                    page_locs.append(loc)
            elif name == "TemplatePage":
                loc = (el.get("BaseLoc") or elem_text(el) or "").strip().lstrip("/")
                if loc and loc not in tpl_locs:
                    tpl_locs.append(loc)

        # 公共资源（字体等）不作为文档块，仅用于定位内容文件
        content_locs: list[str] = []
        for loc in page_locs:
            if loc not in content_locs:
                content_locs.append(loc)
        for loc in tpl_locs:
            if loc not in content_locs:
                content_locs.append(loc)

        def resolve(loc: str) -> str | None:
            if loc in names:
                return loc
            norm = loc.replace("\\", "/")
            for n in names:
                if n.replace("\\", "/").endswith("/" + norm):
                    return n
            return None

        for loc in content_locs:
            resolved = resolve(loc)
            if not resolved:
                continue
            page_xml = ET.fromstring(zf.read(resolved))
            page_no: int | None = None
            for part in resolved.split("/"):
                if not part.lower().startswith("page"):
                    continue
                for seg in part.split("_"):
                    if seg.isdigit():
                        page_no = int(seg)
                        break
                if page_no is not None:
                    break
            texts: list[str] = []
            for el in page_xml.iter():
                if local(el.tag) != "TextCode":
                    continue
                text = elem_text(el)
                if text:
                    texts.append(text)
            if texts:
                blocks.append(
                    {
                        "block_type": "paragraph",
                        "page_no": page_no,
                        "block_index": idx,
                        "text_content": "\n".join(texts),
                        "extra": {"source": "ofd-text-layer", "page": page_no},
                    }
                )
                idx += 1

    if not blocks:
        raise ValueError("OFD 无文本层（扫描版请走视觉模型）")
    return blocks
