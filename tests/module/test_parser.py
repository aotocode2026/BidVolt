from __future__ import annotations

import zipfile
from pathlib import Path

from app.services import parser


def test_parse_txt(tmp_path):
    p = tmp_path / "a.txt"
    p.write_text("第一行\n第二行", encoding="utf-8")
    blocks = parser.parse_to_blocks(p, ".txt")
    assert len(blocks) == 1
    assert "第一行" in blocks[0]["text_content"]


def test_parse_txt_gbk_fallback(tmp_path):
    """中文常见 GBK 编码文本：utf-8 失败回退 gbk，不得产出乱码块。"""
    p = tmp_path / "a.csv"
    p.write_bytes("名称,金额\n电缆,100\n".encode("gbk"))
    blocks = parser.parse_to_blocks(p, ".csv")
    assert "电缆" in blocks[0]["text_content"]


def test_parse_docx(tmp_path):
    from docx import Document

    p = tmp_path / "a.docx"
    doc = Document()
    doc.add_paragraph("商务响应测试")
    doc.add_paragraph("技术响应测试")
    doc.save(str(p))
    blocks = parser.parse_to_blocks(p, ".docx")
    texts = [b["text_content"] for b in blocks]
    assert "商务响应测试" in texts
    assert "技术响应测试" in texts


def test_parse_legacy_doc_requires_libreoffice(tmp_path, monkeypatch):
    """无 LibreOffice 时 .doc 给出可操作提示，而不是"不支持的格式"。"""
    import shutil as _shutil

    monkeypatch.setattr(_shutil, "which", lambda name: None)
    p = tmp_path / "a.doc"
    p.write_bytes(b"\xd0\xcf\x11\xe0")
    import pytest

    with pytest.raises(ValueError) as exc:
        parser.parse_to_blocks(p, ".doc")
    assert "另存为" in str(exc.value)


def test_parse_legacy_doc_roundtrip_docx_keeps_tables(tmp_path):
    """有 LibreOffice（CI 容器）时：.doc 转 docx 解析，表格结构保留。"""
    import shutil
    import subprocess

    import pytest
    from docx import Document

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice is None:
        pytest.skip("本机无 LibreOffice，.doc 转换路径在 CI/容器验证")

    docx_p = tmp_path / "a.docx"
    doc = Document()
    doc.add_paragraph("技术规范书测试")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "序号"
    table.rows[0].cells[1].text = "服务名称"
    table.rows[1].cells[0].text = "1"
    table.rows[1].cells[1].text = "接口调试"
    doc.save(str(docx_p))

    doc_p = tmp_path / "a.doc"
    subprocess.run(
        [soffice, "--headless", "--norestore", "--convert-to", "doc", "--outdir", str(tmp_path), str(docx_p)],
        capture_output=True,
        timeout=180,
        check=True,
    )
    assert doc_p.exists()

    blocks = parser.parse_to_blocks(doc_p, ".doc")
    texts = [b["text_content"] for b in blocks]
    assert any("技术规范书测试" in t for t in texts)
    assert any(b["block_type"] == "table" and "接口调试" in b["text_content"] for b in blocks)
    assert all(b.get("extra", {}).get("source") == "libreoffice-doc" for b in blocks)


def test_parse_docx_header_footer(tmp_path):
    """docx 页眉/页脚进入索引块（正文解析不覆盖的区域）。"""
    from docx import Document

    p = tmp_path / "a.docx"
    doc = Document()
    doc.add_paragraph("正文内容")
    doc.sections[0].header.paragraphs[0].text = "页眉公司名"
    doc.sections[0].footer.paragraphs[0].text = "第 X 页"
    doc.save(str(p))
    blocks = parser.parse_to_blocks(p, ".docx")
    hf = [b for b in blocks if b["block_type"] == "header_footer"]
    assert any("页眉公司名" in b["text_content"] for b in hf)
    assert any("第 X 页" in b["text_content"] for b in hf)


def test_parse_legacy_xls_requires_libreoffice(tmp_path, monkeypatch):
    import shutil as _shutil

    monkeypatch.setattr(_shutil, "which", lambda name: None)
    p = tmp_path / "a.xls"
    p.write_bytes(b"\xd0\xcf\x11\xe0")
    import pytest

    with pytest.raises(ValueError) as exc:
        parser.parse_to_blocks(p, ".xls")
    assert "另存为" in str(exc.value)


def test_parse_legacy_ppt_requires_libreoffice(tmp_path, monkeypatch):
    import shutil as _shutil

    monkeypatch.setattr(_shutil, "which", lambda name: None)
    p = tmp_path / "a.ppt"
    p.write_bytes(b"\xd0\xcf\x11\xe0")
    import pytest

    with pytest.raises(ValueError) as exc:
        parser.parse_to_blocks(p, ".ppt")
    assert "另存为" in str(exc.value)


def _soffice_or_skip():
    import shutil

    import pytest

    s = shutil.which("soffice") or shutil.which("libreoffice")
    if s is None:
        pytest.skip("本机无 LibreOffice，旧格式转换路径在 CI/容器验证")
    return s


def test_parse_legacy_xls_roundtrip_keeps_sheets(tmp_path):
    """.xls → LibreOffice → xlsx → 表格行保留。"""
    import subprocess

    from openpyxl import Workbook

    soffice = _soffice_or_skip()
    xlsx_p = tmp_path / "a.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.append(["材料", "单价"])
    ws.append(["电缆", "120"])
    wb.save(str(xlsx_p))

    xls_p = tmp_path / "a.xls"
    subprocess.run(
        [soffice, "--headless", "--norestore", "--convert-to", "xls", "--outdir", str(tmp_path), str(xlsx_p)],
        capture_output=True,
        timeout=180,
        check=True,
    )
    assert xls_p.exists()

    blocks = parser.parse_to_blocks(xls_p, ".xls")
    assert any(b["block_type"] == "table" and "电缆" in b["text_content"] for b in blocks)
    assert all(b.get("extra", {}).get("source") == "libreoffice-xls" for b in blocks)


def test_parse_legacy_ppt_roundtrip(tmp_path):
    """.ppt → LibreOffice → pptx → 每页文本保留。"""
    import subprocess

    pytest = __import__("pytest")
    pytest.importorskip("pptx")
    from pptx import Presentation

    soffice = _soffice_or_skip()
    pptx_p = tmp_path / "a.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "技术方案汇报"
    prs.save(str(pptx_p))

    ppt_p = tmp_path / "a.ppt"
    subprocess.run(
        [soffice, "--headless", "--norestore", "--convert-to", "ppt", "--outdir", str(tmp_path), str(pptx_p)],
        capture_output=True,
        timeout=180,
        check=True,
    )
    assert ppt_p.exists()

    blocks = parser.parse_to_blocks(ppt_p, ".ppt")
    assert any("技术方案汇报" in b["text_content"] for b in blocks)
    assert all(b.get("extra", {}).get("source") == "libreoffice-ppt" for b in blocks)


def test_parse_xlsx(tmp_path):
    from openpyxl import Workbook

    p = tmp_path / "a.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.append(["材料", "数量"])
    ws.append(["电缆", "100"])
    wb.save(str(p))
    blocks = parser.parse_to_blocks(p, ".xlsx")
    assert any("电缆" in (b.get("text_content") or "") for b in blocks)


def test_parse_xlsx_extensionless_storage_path(tmp_path):
    """Issue #13 输入矩阵暴露：生产存储对象路径为内容寻址的 ".../<sha>/original"
    （无扩展名），openpyxl 按文件名扩展名校验会拒绝该路径
    （InvalidFileException: "does not support  file format"）。
    本地测试此前用 a.xlsx 带扩展名临时文件，从未覆盖此差异。"""
    from openpyxl import Workbook

    p = tmp_path / "a.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.append(["材料", "数量"])
    ws.append(["电缆", "100"])
    wb.save(str(p))

    original = tmp_path / "original"  # 模拟 storage.save 的无扩展名对象路径
    original.write_bytes(p.read_bytes())
    blocks = parser.parse_to_blocks(original, ".xlsx")
    assert any("电缆" in (b.get("text_content") or "") for b in blocks)


def test_parse_pdf(tmp_path):
    try:
        import fitz  # pymupdf
    except ImportError:
        import pytest

        pytest.skip("pymupdf 未安装")
    p = tmp_path / "a.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "招标技术要求：电压等级 10kV")
    doc.save(str(p))
    doc.close()
    blocks = parser.parse_to_blocks(p, ".pdf")
    assert any("10kV" in (b.get("text_content") or "") for b in blocks)


def _write_ofd(
    tmp_path: Path,
    text: str,
    size_attr: str = "Size",
    docroot_loc: str = "Doc_0/Document.xml",
    page_base_loc: str = "Doc_0/Pages/Page_1/Content.xml",
) -> Path:
    """构造最小 OFD（zip+XML），按 TextObject 属性大小写区分标准/easyofd 风格。"""
    p = tmp_path / "a.ofd"
    ns = "http://www.ofdspec.org/2016"
    ofd_xml = f'''<?xml version="1.0" encoding="UTF-8"?>
<ofd:OFD xmlns:ofd="{ns}">
  <ofd:DocBody>
    <ofd:DocInfo><ofd:DocID>12345678901234567890123456789012</ofd:DocID></ofd:DocInfo>
    <ofd:DocRoot><ofd:BaseLoc>{docroot_loc}</ofd:BaseLoc></ofd:DocRoot>
  </ofd:DocBody>
</ofd:OFD>
'''
    document_xml = f'''<?xml version="1.0" encoding="UTF-8"?>
<ofd:Document xmlns:ofd="{ns}">
  <ofd:CommonData>
    <ofd:PageArea><ofd:PhysicalBox>0 0 595.0 842.0</ofd:PhysicalBox></ofd:PageArea>
    <ofd:PublicRes>Doc_0/PublicRes.xml</ofd:PublicRes>
    <ofd:DocumentRes>Doc_0/DocumentRes.xml</ofd:DocumentRes>
  </ofd:CommonData>
  <ofd:Pages>
    <ofd:Page ID="1"><ofd:BaseLoc>{page_base_loc}</ofd:BaseLoc></ofd:Page>
  </ofd:Pages>
</ofd:Document>
'''
    content_xml = f'''<?xml version="1.0" encoding="UTF-8"?>
<ofd:Page xmlns:ofd="{ns}">
  <ofd:Content>
    <ofd:Layer ID="4" Type="Foreground">
      <ofd:TextObject ID="5" Font="3" {size_attr}="12" Boundary="0 0 100 20">
        <ofd:TextCode X="0" Y="12">{text}</ofd:TextCode>
      </ofd:TextObject>
    </ofd:Layer>
  </ofd:Content>
</ofd:Page>
'''
    publicres_xml = f'''<?xml version="1.0" encoding="UTF-8"?>
<ofd:Res xmlns:ofd="{ns}">
  <ofd:Fonts><ofd:Font ID="3" FontName="SimSun" FamilyName="SimSun"/></ofd:Fonts>
</ofd:Res>
'''
    documentres_xml = f'''<?xml version="1.0" encoding="UTF-8"?>
<ofd:Res xmlns:ofd="{ns}"><ofd:MultiMedias/></ofd:Res>
'''
    files = {
        "OFD.xml": ofd_xml,
        "Doc_0/Document.xml": document_xml,
        "Doc_0/Pages/Page_1/Content.xml": content_xml,
        "Doc_0/PublicRes.xml": publicres_xml,
        "Doc_0/DocumentRes.xml": documentres_xml,
    }
    with zipfile.ZipFile(p, "w", zipfile.ZIP_DEFLATED) as zf:
        for name, data in files.items():
            zf.writestr(name, data)
    return p


def test_parse_ofd_standard(tmp_path):
    p = _write_ofd(tmp_path, "招标技术要求：电压等级 10kV", size_attr="Size")
    blocks = parser.parse_to_blocks(p, ".ofd")
    assert len(blocks) == 1
    assert blocks[0]["page_no"] == 1
    assert "10kV" in blocks[0]["text_content"]


def test_parse_ofd_easyofd_lowercase_size(tmp_path):
    """easyofd 生成的文件用小写 size 属性，解析端需兼容。"""
    p = _write_ofd(tmp_path, "商务响应：交付周期 30 天", size_attr="size")
    blocks = parser.parse_to_blocks(p, ".ofd")
    assert any("交付周期" in b.get("text_content") for b in blocks)


def test_parse_ofd_no_text_layer(tmp_path):
    import pytest

    p = tmp_path / "scan.ofd"
    ns = "http://www.ofdspec.org/2016"
    with zipfile.ZipFile(p, "w") as zf:
        zf.writestr(
            "OFD.xml",
            f'<?xml version="1.0" encoding="UTF-8"?><ofd:OFD xmlns:ofd="{ns}"><ofd:DocBody>'
            f'<ofd:DocRoot><ofd:BaseLoc>Doc_0/Document.xml</ofd:BaseLoc></ofd:DocRoot>'
            f"</ofd:DocBody></ofd:OFD>",
        )
        zf.writestr(
            "Doc_0/Document.xml",
            f'<?xml version="1.0" encoding="UTF-8"?><ofd:Document xmlns:ofd="{ns}"><ofd:Pages>'
            f'<ofd:Page ID="1"><ofd:BaseLoc>Doc_0/Pages/Page_0/Content.xml</ofd:BaseLoc></ofd:Page>'
            f"</ofd:Pages></ofd:Document>",
        )
        zf.writestr(
            "Doc_0/Pages/Page_0/Content.xml",
            f'<?xml version="1.0" encoding="UTF-8"?><ofd:Page xmlns:ofd="{ns}"><ofd:Content><ofd:Layer ID="4"/>'
            f"</ofd:Content></ofd:Page>",
        )
    with pytest.raises(ValueError, match="无文本层"):
        parser.parse_to_blocks(p, ".ofd")


def test_parse_ofd_leading_slash_loc(tmp_path):
    """真实 OFD（iOFD/dms360 生成）DocRoot 与 Page BaseLoc 使用 / 前缀绝对路径，需兼容。"""
    p = _write_ofd(
        tmp_path,
        "带签章的 OFD：投标文件",
        docroot_loc="/Doc_0/Document.xml",
        page_base_loc="/Doc_0/Pages/Page_1/Content.xml",
    )
    blocks = parser.parse_to_blocks(p, ".ofd")
    assert len(blocks) == 1
    assert "投标文件" in blocks[0]["text_content"]
    assert blocks[0]["page_no"] == 1


def _write_minimal_pptx(tmp_path: Path, slides: list[str]) -> Path:
    """构造最小 pptx（zip + slideN.xml，a:t 文本），供标准库提取路径测试。"""
    p = tmp_path / "a.pptx"
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    with zipfile.ZipFile(p, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(
            "[Content_Types].xml",
            '<?xml version="1.0"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"/>',
        )
        for i, text in enumerate(slides, start=1):
            zf.writestr(
                f"ppt/slides/slide{i}.xml",
                f'<?xml version="1.0"?><p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
                f'xmlns:a="{ns}"><p:sp><p:txBody><a:p><a:r><a:t>{text}</a:t></a:r></a:p></p:txBody></p:sp></p:sld>',
            )
    return p


def test_parse_pptx_stdlib(tmp_path):
    """Issue #13：供应商注意事项 .pptx 曾报"不支持的格式"——标准库 zip+XML 提取 slide 文本。"""
    p = _write_minimal_pptx(tmp_path, ["投标注意事项：保证金 5%", "第二页内容：交付周期 30 天"])
    blocks = parser._parse_pptx_stdlib(p)
    assert len(blocks) == 2
    assert blocks[0]["page_no"] == 1
    assert blocks[1]["page_no"] == 2
    assert "保证金 5%" in blocks[0]["text_content"]
    assert "交付周期 30 天" in blocks[1]["text_content"]


def test_parse_pptx_no_text_layer(tmp_path):
    import pytest

    p = tmp_path / "empty.pptx"
    with zipfile.ZipFile(p, "w") as zf:
        zf.writestr(
            "ppt/slides/slide1.xml",
            '<?xml version="1.0"?><p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:sp/></p:sld>',
        )
    with pytest.raises(ValueError, match="无文本层"):
        parser._parse_pptx_stdlib(p)
